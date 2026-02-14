from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Function to add a slide with bullet points
    def add_bullet_slide(title_text, bullets):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AnalystSage AI: Privacy-First Document Intelligence"
    subtitle.text = "Revolutionizing Contract Auditing with Retrieval-Augmented Generation (RAG)"

    # Slide 2: The Knowledge Gap
    add_bullet_slide("The Problem with Standard AI", [
        "Static Knowledge: LLMs are frozen in time (training cutoffs).",
        "Hallucinations: AI guesses when it doesn't know your specific data.",
        "Data Privacy: Sending sensitive contracts to the cloud is a risk."
    ])

    # Slide 3: What is RAG?
    add_bullet_slide("What is Retrieval-Augmented Generation (RAG)?", [
        "A bridge between the LLM's reasoning and your private data.",
        "Concept: 'Giving the AI an Open-Book Exam.'",
        "Retrieval: Fact-finding from your documents.",
        "Augmentation: Adding facts to the AI's prompt.",
        "Generation: Producing accurate, cited answers."
    ])

    # Slide 4: The 3-Step Process
    add_bullet_slide("How RAG Works Under the Hood", [
        "Step 1: Embed - Documents turned into vectors and stored in ChromaDB.",
        "Step 2: Search - Query pulls the most relevant 'chunks' of text.",
        "Step 3: Answer - Ollama (Mistral) reads chunks and answers locally."
    ])

    # Slide 5: Introducing AnalystSage AI
    add_bullet_slide("Introducing AnalystSage AI", [
        "Professional Document Intelligence Hub",
        "Semantic Chat: Natural conversation with your library.",
        "Smart Auditor: Automated compliance checking.",
        "Local Core: 100% private, no cloud required."
    ])

    # Slide 6: Use Case: Smart Contract Auditing
    add_bullet_slide("Solving the 'Legal Bottleneck'", [
        "Index Standard: Upload your 'Golden Source' (Approved terms).",
        "Upload Draft: Drop in the new vendor's contract.",
        "Cross-Audit: AI compares terms against the baseline automatically."
    ])

    # Slide 7: The Auditor Scorecard
    add_bullet_slide("Data-Driven Legal Decisions", [
        "KPI Scoring: Liability, IP, and Financial risk percentages.",
        "Deviation Flagging: Red-lining high-risk differences.",
        "Confidence Metrics: Built-in AI certainty tracking."
    ])

    # Slide 8: Technical Excellence
    add_bullet_slide("The Tech Stack", [
        "Models: Ollama (Mistral 7B) + Nomic Embeddings.",
        "Database: ChromaDB (Vector Store).",
        "UI: Modern Streamlit with Glassmorphism design.",
        "Processing: 100% Local (CPU/GPU acceleration)."
    ])

    # Slide 9: Why it Wins
    add_bullet_slide("Why Local RAG?", [
        "Security: Contracts stay on your local disk.",
        "Speed: No network latency or API rate limits.",
        "Cost: Zero token costs or monthly subscriptions."
    ])

    # Slide 10: Conclusion
    add_bullet_slide("The Future of Document Intelligence", [
        "Summary: AnalystSage AI turns hundreds of pages into instant, actionable intelligence.",
        "Closing: 'Intelligence you can trust. Data you can keep.'"
    ])

    # Save presentation
    output_path = "c:\\Automations\\BABOK RAG\\AnalystSage_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
